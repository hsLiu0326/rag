"""Multi-format document loaders.

Supported formats: PDF (via MinerU GPU), DOCX, Markdown, TXT.
PDF → MinerU (layout-aware structured Markdown with tables/formulas)
DOCX → python-docx (paragraphs, tables, heading styles)
MD   → direct read (preserves structure)
TXT  → encoding-detected read
"""

from __future__ import annotations

import io
import re
from pathlib import Path
from typing import Any

import chardet


# ── Shared document type ────────────────────────────────────────────

class Document:
    """Simple document object."""
    page_content: str
    metadata: dict[str, Any]

    def __init__(self, page_content: str = "", metadata: dict[str, Any] | None = None) -> None:
        self.page_content = page_content
        self.metadata = metadata or {}

    def __repr__(self) -> str:
        src = self.metadata.get("source", "?")
        return f"Document(source={src!r}, len={len(self.page_content)})"


# ── PDF Loader ──────────────────────────────────────────────────────

class PDFLoader:
    """Smart PDF loader: fast path for text PDFs, MinerU for complex/scanned.

    - Text PDFs (most docs): pypdf extracts text in <5s, no GPU needed.
    - Scanned/image-heavy PDFs: falls back to MinerU GPU pipeline (~70s cold start).
    """

    def load(self, file_path: str | Path) -> list[Document]:
        file_path = Path(file_path)

        # Fast path: try pypdf first
        from pypdf import PdfReader
        reader = PdfReader(str(file_path))
        pages_text: list[str] = []
        for page in reader.pages:
            t = (page.extract_text() or "").strip()
            if t:
                pages_text.append(t)

        total_chars = sum(len(t) for t in pages_text)

        # If we got meaningful text, use it directly (fast, <1s)
        if total_chars > 50:
            content = "\n\n".join(pages_text)
            print(f"[loader] PDF fast path: {file_path.name} ({len(reader.pages)}p, {total_chars} chars)")
            return [Document(
                page_content=content,
                metadata={"source": file_path.name, "format": "pdf", "parser": "pypdf"},
            )]

        # Fallback: scanned PDF → MinerU GPU
        print(f"[loader] PDF slow path (MinerU): {file_path.name} ({len(reader.pages)}p, only {total_chars} chars text)")
        from src.core.preprocessor.mineru_parser import MinerUParser
        parser = MinerUParser(backend="pipeline")
        markdown_text = parser.parse(str(file_path))
        return [Document(
            page_content=markdown_text,
            metadata={"source": file_path.name, "format": "pdf", "parser": "mineru"},
        )]


# ── DOCX Loader ─────────────────────────────────────────────────────

class DocxLoader:
    """Load text from DOCX files using python-docx, preserving heading styles."""

    def load(self, file_path: str | Path) -> list[Document]:
        from docx import Document as DocxDocument

        file_path = Path(file_path)
        doc = DocxDocument(str(file_path))
        full_text: list[str] = []

        for para in doc.paragraphs:
            style_name = para.style.name if para.style else ""
            text = para.text.strip()
            if not text:
                full_text.append("")
                continue

            if style_name.startswith("Heading"):
                try:
                    level = int(style_name.split()[-1])
                except (ValueError, IndexError):
                    level = 2
                full_text.append(f"{'#' * level} {text}")
            else:
                full_text.append(text)

        # Extract tables as Markdown
        for table in doc.tables:
            rows: list[str] = []
            for row in table.rows:
                cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                rows.append("| " + " | ".join(cells) + " |")
            if rows:
                sep = "|" + "|".join([" --- " for _ in rows[0].split("|")[1:-1]]) + "|"
                rows.insert(1, sep)
                full_text.append("")
                full_text.append("\n".join(rows))
                full_text.append("")

        content = "\n".join(full_text)
        return [Document(
            page_content=content,
            metadata={"source": file_path.name, "format": "docx"},
        )]


# ── Markdown Loader ─────────────────────────────────────────────────

class MarkdownLoader:
    """Load Markdown files, preserving structure."""

    def load(self, file_path: str | Path) -> list[Document]:
        file_path = Path(file_path)
        content = file_path.read_text(encoding="utf-8")
        return [Document(
            page_content=content,
            metadata={"source": file_path.name, "format": "markdown"},
        )]


# ── TXT Loader ──────────────────────────────────────────────────────

class TxtLoader:
    """Load plain text files with encoding detection."""

    def load(self, file_path: str | Path) -> list[Document]:
        file_path = Path(file_path)
        raw = file_path.read_bytes()
        detected = chardet.detect(raw)
        encoding = detected.get("encoding", "utf-8") or "utf-8"
        content = raw.decode(encoding, errors="replace")
        return [Document(
            page_content=content,
            metadata={"source": file_path.name, "format": "txt", "encoding": encoding},
        )]


# ── PPTX Loader ──────────────────────────────────────────────────────

class PptxLoader:
    """Load PowerPoint files via python-pptx. Extracts slide text + notes."""

    def load(self, file_path: str | Path) -> list[Document]:
        from pptx import Presentation

        file_path = Path(file_path)
        prs = Presentation(str(file_path))
        slides: list[str] = []

        for i, slide in enumerate(prs.slides, 1):
            parts: list[str] = [f"## 幻灯片 {i}"]
            # Extract title
            if slide.shapes.title and slide.shapes.title.text.strip():
                parts.append(f"### {slide.shapes.title.text.strip()}")
            # Extract text from all shapes
            for shape in slide.shapes:
                if shape.has_text_frame and shape != slide.shapes.title:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            parts.append(t)
            # Notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    parts.append(f"> {notes}")
            slides.append("\n".join(parts))

        content = "\n\n".join(slides)
        return [Document(
            page_content=content,
            metadata={"source": file_path.name, "format": "pptx"},
        )]


# ── XLSX Loader ──────────────────────────────────────────────────────

class XlsxLoader:
    """Load Excel files via openpyxl. Each sheet becomes a markdown table."""

    def load(self, file_path: str | Path) -> list[Document]:
        from openpyxl import load_workbook

        file_path = Path(file_path)
        wb = load_workbook(str(file_path), read_only=True, data_only=True)
        sheets: list[str] = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = list(ws.iter_rows(max_row=min(ws.max_row or 0, 500), values_only=True))
            if not rows:
                continue

            parts: list[str] = [f"## {sheet_name}"]
            # Build markdown table
            table_lines: list[str] = []
            for r_idx, row in enumerate(rows):
                cells = [str(c) if c is not None else "" for c in row]
                table_lines.append("| " + " | ".join(cells) + " |")
                if r_idx == 0:
                    sep = "|" + "|".join([" --- " for _ in cells]) + "|"
                    table_lines.append(sep)

            parts.extend(table_lines)
            sheets.append("\n".join(parts))

        wb.close()
        content = "\n\n".join(sheets)
        return [Document(
            page_content=content,
            metadata={"source": file_path.name, "format": "xlsx"},
        )]


# ── Loader Registry ─────────────────────────────────────────────────

LOADER_MAP: dict[str, type] = {
    ".pdf": PDFLoader,
    ".docx": DocxLoader,
    ".doc": DocxLoader,
    ".pptx": PptxLoader,
    ".ppt": PptxLoader,
    ".xlsx": XlsxLoader,
    ".xls": XlsxLoader,
    ".md": MarkdownLoader,
    ".txt": TxtLoader,
    ".text": TxtLoader,
}


def load_document(file_path: str | Path) -> list[Document]:
    """Auto-detect format and load document."""
    file_path = Path(file_path)
    suffix = file_path.suffix.lower()
    loader_cls = LOADER_MAP.get(suffix)
    if loader_cls is None:
        raise ValueError(f"Unsupported file format: {suffix}")
    return loader_cls().load(file_path)
