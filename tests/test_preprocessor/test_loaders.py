"""Tests for multi-format document loaders."""

import pytest

from src.core.preprocessor.loaders import (
    LOADER_MAP,
    DocxLoader,
    MarkdownLoader,
    PptxLoader,
    TxtLoader,
    XlsxLoader,
    load_document,
)


def test_loader_registry_covers_all_supported_extensions():
    for ext in (".pdf", ".docx", ".doc", ".pptx", ".ppt", ".xlsx", ".xls", ".md", ".txt", ".text"):
        assert ext in LOADER_MAP


def test_markdown_loader_preserves_content(tmp_path):
    path = tmp_path / "note.md"
    path.write_text("# 标题\n\n正文内容", encoding="utf-8")
    docs = MarkdownLoader().load(path)
    assert len(docs) == 1
    assert docs[0].page_content == "# 标题\n\n正文内容"
    assert docs[0].metadata["source"] == "note.md"


def test_txt_loader_detects_gbk_encoding(tmp_path):
    path = tmp_path / "gbk.txt"
    path.write_bytes("你好，世界，这是中文测试。".encode("gbk"))
    docs = TxtLoader().load(path)
    assert "你好，世界" in docs[0].page_content
    assert docs[0].metadata["encoding"].lower().startswith("gb")


def test_docx_loader_preserves_headings_and_tables(tmp_path):
    from docx import Document

    doc = Document()
    doc.add_heading("标题一", level=1)
    doc.add_paragraph("普通段落内容")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "单元格A"
    table.cell(1, 1).text = "单元格B"
    path = tmp_path / "sample.docx"
    doc.save(str(path))

    docs = DocxLoader().load(path)
    content = docs[0].page_content
    assert "# 标题一" in content
    assert "普通段落内容" in content
    assert "单元格A" in content
    assert "| 单元格A |" in content


def test_pptx_loader_extracts_slide_text(tmp_path):
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "测试标题"
    slide.placeholders[1].text = "正文内容"
    path = tmp_path / "slides.pptx"
    prs.save(str(path))

    docs = PptxLoader().load(path)
    content = docs[0].page_content
    assert "## 幻灯片 1" in content
    assert "### 测试标题" in content
    assert "正文内容" in content


def test_xlsx_loader_converts_sheets_to_markdown_tables(tmp_path):
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "成绩表"
    ws.append(["姓名", "分数"])
    ws.append(["小明", 95])
    path = tmp_path / "grades.xlsx"
    wb.save(str(path))

    docs = XlsxLoader().load(path)
    content = docs[0].page_content
    assert "## 成绩表" in content
    assert "| 姓名 | 分数 |" in content
    assert "| 小明 | 95 |" in content


def test_load_document_dispatches_by_extension(tmp_path):
    md = tmp_path / "a.md"
    md.write_text("正文", encoding="utf-8")
    assert load_document(md)[0].metadata["format"] == "markdown"

    txt = tmp_path / "b.txt"
    txt.write_text("正文", encoding="utf-8")
    assert load_document(txt)[0].metadata["format"] == "txt"


def test_load_document_rejects_unknown_format(tmp_path):
    unknown = tmp_path / "file.xyz"
    unknown.write_text("x", encoding="utf-8")
    with pytest.raises(ValueError, match="Unsupported file format"):
        load_document(unknown)
